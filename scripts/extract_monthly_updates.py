# -*- coding: utf-8 -*-
"""
Extract Monthly Solution Updates from PowerPoint Presentations
===============================================================
Parses .pptx files from Monthly Project Status Updates folder.
Maps solution names to core_ids using MO-DB_Solutions database.
Outputs an Excel file (.xlsx) for review before import to MO-DB_Updates.

Usage: python extract_monthly_updates.py
"""

from pptx import Presentation
from pathlib import Path
import pandas as pd
import re
from datetime import datetime
import uuid
import sys

# Configuration
BASE_PATH = Path(r"C:\Users\cjtucke3\Documents\Personal\MO-development\source-archives\Monthly Project Status Updates")
SOLUTIONS_PATH = Path(r"C:\Users\cjtucke3\Documents\Personal\MO-development\nsite-mo-viewer\database-files\MO-Viewer Databases\MO-DB_Solutions.xlsx")
OUTPUT_PATH = Path(r"C:\Users\cjtucke3\Documents\Personal\MO-development\nsite-mo-viewer\database-files\monthly_updates_import.xlsx")

# CSV columns matching MO-DB_Updates
OUTPUT_COLUMNS = [
    'update_id', 'solution_id', 'update_text', 'source_document',
    'source_category', 'source_url', 'source_tab', 'meeting_date',
    'created_at', 'created_by', 'slide_title', 'presenter'
]

# Known solution name patterns and mappings
# These help with matching slide titles to core_ids
SOLUTION_ALIASES = {
    'harmonized landsat sentinel': 'HLS',
    'harmonized landsat sentinel 2': 'HLS',
    'hls': 'HLS',
    'opera': 'OPERA',
    'dswx': 'DSWx',
    'dist': 'DIST',
    'disp': 'DISP',
    'icesat-2': 'ICESat-2',
    'icesat2': 'ICESat-2',
    'rtc-s1': 'RTC-S1',
    'tempo nrt': 'TEMPO NRT',
    'tempo validation': 'TEMPO Validation',
    'tempo nrt enhanced': 'TEMPO-NRT-Enhanced',
    'tempo enhanced': 'TEMPO-NRT-Enhanced',
    'pbl': 'PBL',
    'vlm': 'VLM',
    'vertical land motion': 'VLM',
    'gaban': 'GABAN',
    'global algal bloom': 'GABAN',
    'hls low latency': 'HLS-LL',
    'hls ll': 'HLS-LL',
    'hls-vi': 'HLS-VI',
    'admg': 'ADMG',
    'airborne data management': 'ADMG',
    'esdis': 'ESDIS',
    'lance': 'LANCE',
    'firms': 'FIRMS',
    'internet of animals': 'IoA',
    'ioa': 'IoA',
    'csda': 'CSDA',
    'air quality': 'Air Quality',
    'air quality: pandora sensors': 'AQ-Pandora',
    'air quality: pm2.5': 'AQ-PM2.5',
    'air quality: gmao': 'AQ-GMAO',
    'pm2.5': 'AQ-PM2.5',
    'pandora': 'AQ-Pandora',
    'pandora sensors': 'AQ-Pandora',
    'gmao': 'AQ-GMAO',
    'gacr': 'GACR',
    'mwow': 'MWoW',
    'snwg mo': 'SNWG-MO',
    'mo': 'SNWG-MO',
    'management office': 'SNWG-MO',
    'snwg management office': 'SNWG-MO',
    'nsite mo': 'SNWG-MO',
    # Ensure consistent casing
    'aq_pandora': 'AQ-Pandora',
    'aq_pm2.5': 'AQ-PM2.5',
    'aq_gmao': 'AQ-GMAO',
    'tempo_enhanced': 'TEMPO-NRT-Enhanced',
    'tempo-nrt-enhanced': 'TEMPO-NRT-Enhanced',
}

# Slides to skip (not solution-specific updates)
SKIP_TITLE_PATTERNS = [
    r'^agenda',
    r'^snwg\s+(headquarters|hq)',
    r'^solution\s+team\s+folders',
    r'^insert\s+impact',
    r'^cycle\s+\d+\s+\(\d{4}\)',  # "Cycle 1 (2016)" divider slides
    r'^assessment',
    r'^appendix',
    r'^thank\s*you',
    r'^questions',
    r'^backup',
    r'^management\s+office',
    r'^nsite\s+mo\s+slide',
    r'^satellite\s+needs\s+working\s+group\s+management',  # Title slide
    r'^\[project\s+name\]',  # Template placeholder
    r'^project\s+name',  # Template placeholder
    r'^cycle\s+\d+[:\s]',  # "Cycle 2: Land Surface" divider slides
    r'^\[project\]',  # Template placeholder
    r'^\d{4}\s+snwg\s+assessment',  # Assessment slides
    r'^status\s+and\s+near',  # Non-solution slides
    r'^snwg\s+sep',  # SEP overview slides
    r'^sep\s+',  # SEP slides
]

# Illegal characters for Excel (control characters except tab/newline)
ILLEGAL_CHARS_RE = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f]')


def generate_update_id():
    """Generate a unique update ID"""
    date_str = datetime.now().strftime('%Y%m%d')
    random_part = str(uuid.uuid4())[:4].upper()
    return f"UPD_{date_str}_{random_part}"


def sanitize_for_excel(text):
    """Remove illegal characters that can't be written to Excel"""
    if not text:
        return text
    # Remove control characters except newline and tab
    text = ILLEGAL_CHARS_RE.sub('', text)
    # Replace other problematic characters
    text = text.replace('\r', '\n')
    return text


# Post-processing normalization for solution IDs
SOLUTION_ID_NORMALIZATION = {
    'hls': 'HLS',
    'pbl': 'PBL',
    'vlm': 'VLM',
    'admg': 'ADMG',
    'esdis': 'ESDIS',
    'gaban': 'GABAN',
    'lance': 'LANCE',
    'firms': 'FIRMS',
    'csda': 'CSDA',
    'ioa': 'IoA',
    'opera': 'OPERA',
    'dswx': 'DSWx',
    'dist': 'DIST',
    'disp': 'DISP',
    'gacr': 'GACR',
    'mwow': 'MWoW',
    'aq_pandora': 'AQ-Pandora',
    'aq_pm2.5': 'AQ-PM2.5',
    'aq_gmao': 'AQ-GMAO',
    'tempo_enhanced': 'TEMPO-NRT-Enhanced',
    'mo': 'SNWG-MO',
}


def normalize_solution_id(solution_id):
    """Normalize solution ID to canonical form"""
    if not solution_id:
        return solution_id
    lower = solution_id.lower().strip()
    return SOLUTION_ID_NORMALIZATION.get(lower, solution_id)


def build_solution_mapping():
    """Build mapping from solution names/aliases to core_ids"""
    df = pd.read_excel(SOLUTIONS_PATH)
    mapping = {}

    for _, row in df.iterrows():
        core_id = str(row.get('core_id', '')).strip()
        if not core_id or core_id == 'nan':
            continue

        official_name = str(row.get('core_official_name', '')).strip()
        alternates = str(row.get('core_alternate_names', '')).strip()

        # Add official name
        if official_name and official_name != 'nan':
            mapping[official_name.lower()] = core_id
            # Also add without parenthetical
            clean = re.sub(r'\s*\([^)]+\)\s*', '', official_name).strip()
            if clean:
                mapping[clean.lower()] = core_id

        # Add core_id itself as a match
        mapping[core_id.lower()] = core_id

        # Add alternates (pipe or comma separated)
        if alternates and alternates != 'nan':
            for alt in re.split(r'[|,]', alternates):
                alt = alt.strip()
                if alt:
                    mapping[alt.lower()] = core_id

    # Add manual aliases
    for alias, core_id in SOLUTION_ALIASES.items():
        mapping[alias.lower()] = core_id

    return mapping


def find_core_id(text, solution_mapping):
    """Find core_id for a solution name using the mapping"""
    if not text:
        return None

    text_clean = text.strip().lower()

    # Direct match
    if text_clean in solution_mapping:
        return solution_mapping[text_clean]

    # Try without parenthetical suffix
    no_paren = re.sub(r'\s*\([^)]+\)\s*$', '', text_clean).strip()
    if no_paren in solution_mapping:
        return solution_mapping[no_paren]

    # Try partial matches
    for name, core_id in solution_mapping.items():
        if name in text_clean or text_clean in name:
            return core_id

    # Try each word in the text
    words = text_clean.split()
    for word in words:
        if word in solution_mapping:
            return solution_mapping[word]

    return None


def should_skip_slide(title):
    """Check if this slide should be skipped (not a solution update)"""
    if not title:
        return True

    title_lower = title.lower().strip()

    for pattern in SKIP_TITLE_PATTERNS:
        if re.search(pattern, title_lower):
            return True

    return False


def extract_date_from_shape_text(text):
    """Extract date from 'Update as of:' text"""
    if not text:
        return None

    # Pattern: "Update as of: MM/DD/YYYY" or "Update as of: M/D/YY"
    match = re.search(r'update\s+as\s+of[:\s]+(\d{1,2})[/\-](\d{1,2})[/\-](\d{2,4})', text, re.IGNORECASE)
    if match:
        month, day, year = match.groups()
        year = int(year)
        if year < 100:
            year += 2000
        try:
            return f"{year}-{int(month):02d}-{int(day):02d}"
        except:
            pass

    return None


def extract_date_from_filename(filename):
    """Extract meeting date from filename like '2026-01 NSITE Monthly...'"""
    # Pattern: YYYY-MM at start of filename
    match = re.match(r'(\d{4})-(\d{2})', filename)
    if match:
        year, month = match.groups()
        # Use 15th as default day for monthly updates
        return f"{year}-{month}-15"

    # Pattern: YYYY-MM-DD
    match = re.match(r'(\d{4}-\d{2}-\d{2})', filename)
    if match:
        return match.group(1)

    return None


def parse_title(title_text):
    """Parse slide title to extract solution name and presenter"""
    if not title_text:
        return None, None

    # Common patterns:
    # "Solution Name\nPresenter Name, Title"
    # "Solution Name\nPresenter Name"
    # "Solution NamePresenter Name" (no newline)

    lines = title_text.strip().split('\n')

    if len(lines) >= 2:
        solution = lines[0].strip()
        presenter = lines[1].strip()
    else:
        # Try to split on common patterns
        # Look for patterns like "Name, Title" at the end
        text = lines[0]
        # Match: "Solution NamePerson Name, Role"
        match = re.match(r'^(.+?)\s*([A-Z][a-z]+\s+[A-Z][a-z]+(?:,\s*.+)?)$', text)
        if match:
            solution = match.group(1).strip()
            presenter = match.group(2).strip()
        else:
            solution = text
            presenter = None

    return solution, presenter


def extract_updates_from_slide(slide, solution_mapping, file_date, filename):
    """Extract updates from a single slide"""
    updates = []

    # Get title
    title_text = slide.shapes.title.text if slide.shapes.title else None

    if should_skip_slide(title_text):
        return []

    solution_name, presenter = parse_title(title_text)
    core_id = find_core_id(solution_name, solution_mapping)

    if not core_id:
        # Try just the first word of the title (often the acronym)
        first_word = solution_name.split()[0] if solution_name else ''
        core_id = find_core_id(first_word, solution_mapping)

    if not core_id:
        # Check if this is a template placeholder or non-solution slide
        if solution_name:
            name_lower = solution_name.lower()
            skip_patterns = [
                r'\[project',
                r'\[new\s+solution',
                r'snwg\s+assessment',
                r'status\s+and\s+near',
                r'^sep\s',
                r'snwg\s+sep',
                r'we\s+welcome',
            ]
            for pattern in skip_patterns:
                if re.search(pattern, name_lower):
                    return []
        # Use the solution name as-is for manual review
        core_id = solution_name

    # Find the "Update as of" date
    slide_date = None
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            slide_date = extract_date_from_shape_text(shape.text)
            if slide_date:
                break

    # Fall back to filename date
    meeting_date = slide_date or file_date

    # Extract text from all shapes
    all_text_parts = []

    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue

        text = shape.text.strip()
        if not text:
            continue

        # Skip the title (already captured)
        if title_text and text == title_text:
            continue

        # Skip "Update as of:" lines
        if re.search(r'update\s+as\s+of', text, re.IGNORECASE):
            continue

        # Skip "Project Phase:" standalone lines
        if re.match(r'^project\s+phase\s*:', text, re.IGNORECASE):
            continue

        all_text_parts.append(text)

    # Join all content and clean up
    full_text = '\n\n'.join(all_text_parts)

    # Instead of splitting into sections, consolidate into ONE update per slide
    # Clean up the boilerplate and section headers
    cleaned_text = clean_section_headers(full_text)

    # Skip if only boilerplate remains
    if is_template_text(cleaned_text):
        return []

    if len(cleaned_text.strip()) < 30:
        return []

    updates.append({
        'update_id': generate_update_id(),
        'solution_id': normalize_solution_id(core_id),
        'update_text': sanitize_for_excel(cleaned_text),
        'source_document': 'Monthly Status Meeting',
        'source_category': 'MO',
        'source_url': '',
        'source_tab': filename,
        'meeting_date': meeting_date,
        'created_at': datetime.now().isoformat(),
        'created_by': 'monthly_import',
        'slide_title': sanitize_for_excel(solution_name),
        'presenter': sanitize_for_excel(presenter)
    })

    return updates


def parse_sections(text):
    """Parse text into logical sections based on headers"""
    sections = []

    # Common section headers in the presentations
    section_patterns = [
        r'what\s+programmatic.*?milestones',
        r'what\s+software.*?milestones',
        r'what\s+have\s+i\s+done',
        r'the\s+snwg\s+mo\s+can\s+help',
        r'what\s+is\s+new',
        r'status\s+update',
        r'key\s+accomplishments',
        r'next\s+steps',
        r'challenges',
        r'upcoming\s+activities',
    ]

    # Try to split on section headers
    current_section = None
    current_content = []

    for line in text.split('\n'):
        line_stripped = line.strip()
        if not line_stripped:
            continue

        # Check if this is a section header
        is_header = False
        for pattern in section_patterns:
            if re.search(pattern, line_stripped.lower()):
                # Save previous section
                if current_content:
                    content = '\n'.join(current_content)
                    if len(content.strip()) > 20:
                        sections.append((current_section, content))

                current_section = line_stripped
                current_content = []
                is_header = True
                break

        if not is_header:
            current_content.append(line)

    # Save last section
    if current_content:
        content = '\n'.join(current_content)
        if len(content.strip()) > 20:
            sections.append((current_section, content))

    # If no sections found, return whole text as one section
    if not sections and len(text.strip()) > 20:
        sections.append((None, text))

    return sections


def is_template_text(text):
    """Check if text is just template/boilerplate"""
    template_patterns = [
        r'^project\s+status\s*$',
        r'^project\s+status\s*\n\s*project\s+phase',
        r'^project\s+phase\s*:\s*(operations|implementation|planning|development)',
        r'^solution\s+products?\s*\(acronyms',
        r'^solution\s+products?\s*:',
        r'^what\s+have\s+i\s+done.*?\?\s*$',
        r'^what\s+programmatic.*?\?\s*$',
        r'^what\s+software.*?\?\s*$',
        r'^the\s+snwg\s+mo\s+can\s+help.*?\s*$',
        r'^n/a\s*$',
        r'^none\s*$',
        r'^no\s+update\s*$',
        r'^tbd\s*$',
        r'sneaky\s*slide',
        r'^project\s+description\s*$',
    ]

    text_lower = text.lower().strip()

    for pattern in template_patterns:
        if re.search(pattern, text_lower):
            return True

    # Very short text with no real content
    if len(text_lower) < 20:
        return True

    return False


def clean_section_headers(text):
    """Remove or simplify section header boilerplate"""
    # Remove the long question headers but keep as simple labels
    replacements = [
        (r'\*\*What have I done to ensure that my solution is the right fit for end users and/or is achieving its intended impact\?.*?\*\*\s*', 'User Engagement: '),
        (r'\*\*What programmatic/project timeline milestones have occurred this month\?\*\*\s*', 'Milestones: '),
        (r'\*\*What software/hardware/location/product development milestones have occurred this month\?\*\*\s*', 'Development: '),
        (r'\*\*The SNWG MO can help me with this roadblock or challenge:\*\*\s*', 'Needs Help: '),
        (r'Project Status\s*\n\s*Project Phase:\s*\w+\s*\n?', ''),
        (r'Solution Products?\s*\(acronyms and definitions\):\s*', 'Products: '),
        (r'Project Description\s*\n?', ''),
    ]

    result = text
    for pattern, replacement in replacements:
        result = re.sub(pattern, replacement, result, flags=re.IGNORECASE)

    # Clean up multiple newlines
    result = re.sub(r'\n{3,}', '\n\n', result)
    result = result.strip()

    return result


def process_presentation(pptx_path, solution_mapping):
    """Process a single PowerPoint file and extract all updates"""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  Error opening {pptx_path.name}: {e}")
        return []

    file_date = extract_date_from_filename(pptx_path.name)
    updates = []

    for slide in prs.slides:
        slide_updates = extract_updates_from_slide(
            slide, solution_mapping, file_date, pptx_path.name
        )
        updates.extend(slide_updates)

    return updates


def main():
    all_updates = []
    files_processed = 0

    print("Building solution name to core_id mapping...")
    solution_mapping = build_solution_mapping()
    print(f"  Loaded {len(solution_mapping)} name/alias mappings")
    print()

    print("Extracting monthly updates from PowerPoint presentations...")
    print(f"Base path: {BASE_PATH}")
    print()

    # Process files in root directory
    root_files = list(BASE_PATH.glob("*.pptx"))
    if root_files:
        print(f"Processing root folder ({len(root_files)} files)...")
        for pptx_file in sorted(root_files):
            print(f"  {pptx_file.name}")
            updates = process_presentation(pptx_file, solution_mapping)
            if updates:
                all_updates.extend(updates)
            files_processed += 1
        print(f"  Found {sum(1 for u in all_updates)} updates")
        print()

    # Process FY folders
    for fy_folder in sorted(BASE_PATH.glob("FY*")):
        if fy_folder.is_dir():
            pptx_files = list(fy_folder.glob("*.pptx"))
            if pptx_files:
                print(f"Processing {fy_folder.name} ({len(pptx_files)} files)...")
                fy_updates = 0
                for pptx_file in sorted(pptx_files):
                    updates = process_presentation(pptx_file, solution_mapping)
                    if updates:
                        fy_updates += len(updates)
                        all_updates.extend(updates)
                    files_processed += 1
                print(f"  Found {fy_updates} updates")

    # Skip Biweekly presentations (different format)
    biweekly = BASE_PATH / "Biweekly presentations"
    if biweekly.exists():
        print()
        print("Skipping Biweekly presentations folder (different format)")

    print()
    print("=" * 60)
    print(f"Total files processed: {files_processed}")
    print(f"Total updates found: {len(all_updates)}")

    # Count by solution
    by_solution = {}
    for u in all_updates:
        sol = u['solution_id']
        by_solution[sol] = by_solution.get(sol, 0) + 1

    print(f"Unique solutions: {len(by_solution)}")
    print("\nTop 15 solutions by update count:")
    for sol, count in sorted(by_solution.items(), key=lambda x: -x[1])[:15]:
        print(f"  {sol}: {count}")

    # Count by date
    by_date = {}
    for u in all_updates:
        date = u.get('meeting_date', 'Unknown')
        by_date[date] = by_date.get(date, 0) + 1

    print(f"\nUpdates by date (showing {min(10, len(by_date))} most recent):")
    for date, count in sorted(by_date.items(), reverse=True)[:10]:
        print(f"  {date}: {count}")

    # Write Excel file
    if all_updates:
        OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

        df = pd.DataFrame(all_updates, columns=OUTPUT_COLUMNS)

        # Sort by date descending, then solution
        df['meeting_date'] = pd.to_datetime(df['meeting_date'], errors='coerce')
        df = df.sort_values(['meeting_date', 'solution_id'], ascending=[False, True])
        df['meeting_date'] = df['meeting_date'].dt.strftime('%Y-%m-%d')

        # Write to Excel
        with pd.ExcelWriter(OUTPUT_PATH, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Updates', index=False)

            # Auto-adjust column widths
            worksheet = writer.sheets['Updates']
            for idx, col in enumerate(df.columns):
                max_length = max(
                    df[col].astype(str).map(len).max(),
                    len(col)
                )
                # Limit width for readability
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[chr(65 + idx)].width = adjusted_width

        print(f"\nExcel file written to: {OUTPUT_PATH}")
        print(f"Open to review before importing to MO-DB_Updates")
    else:
        print("\nNo updates found to export.")


if __name__ == '__main__':
    main()
