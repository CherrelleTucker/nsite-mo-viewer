"""
Create MO Viewer Implementation Presentation PowerPoint
Matching the exact style of the SEP presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors - matching SEP presentation exactly
DARK_BLUE = RGBColor(0x1a, 0x23, 0x5c)  # Title/header blue
GREEN = RGBColor(0x2e, 0x7d, 0x32)
BLUE = RGBColor(0x19, 0x76, 0xd2)
ORANGE = RGBColor(0xf5, 0x7c, 0x00)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLACK = RGBColor(0x00, 0x00, 0x00)

def add_title_slide(prs, title_text, subtitle_text, date_text):
    """Add a dark blue title slide (like slide 1 and 16)"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.6))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date_text
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title_text, subtitle_text=None):
    """Add a content slide with dark blue header bar"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue header bar
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = DARK_BLUE
    header_bar.line.fill.background()

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    if subtitle_text:
        p.text = f"{title_text}: {subtitle_text}"
    else:
        p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_bullet_paragraph(text_frame, text, level=0, bold=False, font_size=20):
    """Add a bullet point to text frame"""
    if len(text_frame.paragraphs) == 1 and text_frame.paragraphs[0].text == "":
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = BLACK
    p.space_before = Pt(8)
    return p

def add_rounded_box(slide, left, top, width, height, border_color, title_text, title_color=None):
    """Add a rounded rectangle box with colored border and title"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = border_color
    box.line.width = Pt(3)

    # Title inside box
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = title_color if title_color else border_color

    return box

# ============ SLIDE 1: Title ============
add_title_slide(prs, "MO Viewer", "Information Management Platform", "Implementation-View | January 2026")

# ============ SLIDE 2: The Problem ============
slide2 = add_content_slide(prs, "The Problem", "lots of data, not a lot of ways to quickly access the information it informs")

content = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(6))
tf = content.text_frame
tf.word_wrap = True

add_bullet_paragraph(tf, "MO-wide:", bold=True)
add_bullet_paragraph(tf, "over 9000 interconnected files stored in Google Drive, acting as a \"database\"", level=1, font_size=18)
add_bullet_paragraph(tf, "Jenny, Cherrelle, Slack, emails, Teams, and various meetings acting as database interfaces", level=1, font_size=18)
add_bullet_paragraph(tf, "multiple copies of various files used to create actionable information conveyed in meetings", level=0, font_size=20)

add_bullet_paragraph(tf, "Implementation-specific:", bold=True)
add_bullet_paragraph(tf, "48 solutions across 5 assessment cycles (2016-2024)", level=1, font_size=18)
add_bullet_paragraph(tf, "33 data fields tracked per solution (lifecycle phase, milestones, documents, contracts)", level=1, font_size=18)
add_bullet_paragraph(tf, "4 major milestones per solution (ATP, F2I, ORR, Closeout) with 9 document deliverables each", level=1, font_size=18)
add_bullet_paragraph(tf, "~4,200+ stakeholder contacts linked to solutions", level=1, font_size=18)
add_bullet_paragraph(tf, "Historical updates scattered across 4+ years of meeting notes", level=1, font_size=18)

# ============ SLIDE 3: Early Solution Attempts (Visual) ============
slide3 = add_content_slide(prs, "Early Solution attempts")

# Year labels and placeholder boxes for screenshots
years = ["2022-2023", "2024", "2025"]
for i, year in enumerate(years):
    left = Inches(0.8 + i * 4.2)

    # Placeholder rectangle for screenshot
    placeholder = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(4.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xee, 0xee, 0xee)
    placeholder.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

    # Add placeholder text
    ph_text = slide3.shapes.add_textbox(left + Inches(0.5), Inches(3.5), Inches(2.8), Inches(0.5))
    tf = ph_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[Screenshot]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

    # Year label
    year_label = slide3.shapes.add_textbox(left, Inches(6.2), Inches(3.8), Inches(0.5))
    tf = year_label.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 4: Early Solution Attempts (Table) ============
slide4 = add_content_slide(prs, "Early Solution attempts")

# Column headers
years = ["2022-2023", "2024", "2025", "2026"]
col_width = Inches(3.1)
for i, year in enumerate(years):
    left = Inches(0.4 + i * 3.2)
    year_box = slide4.shapes.add_textbox(left, Inches(1.2), col_width, Inches(0.5))
    tf = year_box.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# Row 1: Storage
storage_texts = [
    "files stored in GWorkspace",
    "files stored in GWorkspace",
    "files stored in GWorkspace",
    "files stored in GWorkspace"
]
for i, text in enumerate(storage_texts):
    left = Inches(0.4 + i * 3.2)
    box = slide4.shapes.add_textbox(left, Inches(1.8), col_width, Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)

# Row 2: Data Processing
processing_texts = [
    "data manually entered, tracked, transformed, and analyzed;\nfiles viewed through Google Site",
    "data manually entered, automatically transformed across files with ~40 Google AppsScripts;\nfiles viewed through direct access",
    "data pulled from GWorkspace file, applied to custom Github templates, Google AppScript accesses Github and displays the data",
    "data pulled from Source of Truth files into Databases."
]
for i, text in enumerate(processing_texts):
    left = Inches(0.4 + i * 3.2)
    box = slide4.shapes.add_textbox(left, Inches(2.7), col_width, Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)

# Row 3: Access Model
access_texts = [
    "view-only;\nlinks to source files\nno edit access",
    "linking network of source files;\nfull edit access (working directly in all files)",
    "view-only, linking of source files;\nadditional interface layer between source data and information transformation;\nheavy maintenance lift",
    "direct interaction with maintained data to transform into information as needed,\nwithout risk to Source of Truth files"
]
for i, text in enumerate(access_texts):
    left = Inches(0.4 + i * 3.2)
    box = slide4.shapes.add_textbox(left, Inches(4.6), col_width, Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)

# ============ SLIDE 5: The Meta-Analysis ============
slide5 = add_content_slide(prs, "The Meta-Analysis")

content = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(6))
tf = content.text_frame
tf.word_wrap = True

add_bullet_paragraph(tf, "Conducted analysis of the Implementation tracking process as a proposed Solution to address documented Needs", font_size=22)
add_bullet_paragraph(tf, "Evaluated how solution data flows across all MO activities from multiple perspectives", level=1, font_size=18)

add_bullet_paragraph(tf, "Identified design concerns with existing tracking methods:", font_size=22)
add_bullet_paragraph(tf, "Solution data scattered across individual Drive folders, QuickLook sheets, and meeting notes", level=1, font_size=18)
add_bullet_paragraph(tf, "No single view of \"where is this solution right now?\"", level=1, font_size=18)
add_bullet_paragraph(tf, "Milestone and document status requires opening multiple files to piece together", level=1, font_size=18)

add_bullet_paragraph(tf, "Discovered: ability to consolidate solution metadata into structured database while maintaining links to source files. This change means much more interconnectivity of data with much less code.", font_size=22)

# Bold the key phrase
p = tf.add_paragraph()
p.text = "Single source of truth for solution status with full audit trail back to original documents."
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.space_before = Pt(16)

# ============ SLIDE 6: How We Work ============
slide6 = add_content_slide(prs, "How we work, in terms of an Information Management System")

# Placeholder for whiteboard image
placeholder = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

ph_text = slide6.shapes.add_textbox(Inches(4), Inches(4), Inches(5), Inches(0.5))
tf = ph_text.text_frame
p = tf.paragraphs[0]
p.text = "[Whiteboard Diagram]"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 7: Context - Why MO Viewer ============
slide7 = add_content_slide(prs, "Context", "Why MO Viewer and not just Implementation Tracker?")

# Three rounded boxes with colored borders
box_data = [
    ("SEP", BLUE, [
        "SEP milestones map directly to Implementation lifecycle phases",
        "Stakeholder engagement aligns with solution readiness",
        "Working sessions prepare stakeholders for transitions"
    ]),
    ("Comms", GREEN, [
        "Solution updates become Comms stories and talking points",
        "Milestone achievements are outreach opportunities",
        "Success stories amplify Implementation wins"
    ]),
    ("Implementation", ORANGE, [
        "Actions from Implementation meetings need tracking and assignment",
        "Meeting notes feed solution update database",
        "Availability affects milestone scheduling"
    ])
]

for i, (title, color, bullets) in enumerate(box_data):
    left = Inches(0.4 + i * 4.3)
    add_rounded_box(slide7, left, Inches(1.3), Inches(4.1), Inches(2.5), color, title)

    # Bullets inside box
    bullet_box = slide7.shapes.add_textbox(left + Inches(0.2), Inches(1.9), Inches(3.7), Inches(1.8))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    for bullet in bullets:
        add_bullet_paragraph(tf, bullet, font_size=12)

# "Implementation is the foundation:" section
foundation_label = slide7.shapes.add_textbox(Inches(0.4), Inches(4.0), Inches(12), Inches(0.5))
tf = foundation_label.text_frame
p = tf.paragraphs[0]
p.text = "Implementation is the foundation:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

foundation_items = [
    "SEP - SEP wraps around the Solution Lifecycle framework; solution data had to be built first because SEP engagement maps directly to solution phases and milestones",
    "Comms - Every solution milestone is a potential story; Implementation provides the \"what happened\" that Comms amplifies",
    "Team - Internal planning meetings generate updates and actions that feed back into Implementation tracking",
    "Reports - Quad charts, milestone reports, and historical updates all pull from Implementation data"
]

foundation_box = slide7.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.5))
tf = foundation_box.text_frame
tf.word_wrap = True
for item in foundation_items:
    p = tf.add_paragraph()
    p.text = "✓ " + item
    p.font.size = Pt(14)
    p.space_before = Pt(4)

# ============ SLIDE 8: Answering Implementation's Questions ============
slide8 = add_content_slide(prs, "Answering Implementation's Questions with MO Viewer")

# Three numbered questions with circles
questions = [
    ("1", "Where Are We?", "What phase is this solution in and what milestones are coming up?"),
    ("2", "What's the Status?", "Are documents complete? What deliverables are pending?"),
    ("3", "What's Happening?", "What are the recent updates and who are the stakeholders?")
]

for i, (num, title, desc) in enumerate(questions):
    top = Inches(1.5 + i * 1.7)

    # Number circle
    circle = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), top, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_BLUE
    circle.line.fill.background()

    # Number in circle
    num_box = slide8.shapes.add_textbox(Inches(0.5), top + Inches(0.1), Inches(0.8), Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Question title
    title_box = slide8.shapes.add_textbox(Inches(1.5), top, Inches(4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Question description
    desc_box = slide8.shapes.add_textbox(Inches(1.5), top + Inches(0.5), Inches(5), Inches(0.6))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(18)

# Placeholder for screenshot on right
placeholder = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.3), Inches(6), Inches(5.5))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

ph_text = slide8.shapes.add_textbox(Inches(8.5), Inches(3.8), Inches(3), Inches(0.5))
tf = ph_text.text_frame
p = tf.paragraphs[0]
p.text = "[Screenshot]"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 9: Q1 - Where Are We? ============
slide9 = add_content_slide(prs, "Q1", "Where Are We in the Lifecycle?")

# Two rounded boxes
add_rounded_box(slide9, Inches(0.4), Inches(1.3), Inches(5.8), Inches(2.5), GREEN, "Solution Picker", GREEN)
add_rounded_box(slide9, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.5), BLUE, "Milestone Timeline", BLUE)

# Bullets for Solution Picker
sp_bullets = slide9.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(5.4), Inches(2))
tf = sp_bullets.text_frame
tf.word_wrap = True
for item in ["Filter by cycle (1-5)", "Filter by lifecycle phase", "Filter by group (HLS, OPERA, etc.)", "Search across all solutions"]:
    add_bullet_paragraph(tf, "• " + item, font_size=16)

# Bullets for Milestone Timeline
mt_bullets = slide9.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.4), Inches(2))
tf = mt_bullets.text_frame
tf.word_wrap = True
for item in ["Visual progress indicators", "ATP, F2I, ORR, Closeout tracking", "Color-coded completion status", "Date-based milestone planning"]:
    add_bullet_paragraph(tf, "• " + item, font_size=16)

# Lifecycle phases section
phases_label = slide9.shapes.add_textbox(Inches(0.4), Inches(4.0), Inches(12), Inches(0.5))
tf = phases_label.text_frame
p = tf.paragraphs[0]
p.text = "Lifecycle Phases Tracked in MO Viewer:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

phases_box = slide9.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.5))
tf = phases_box.text_frame
tf.word_wrap = True
phases = [
    ("Formulation", "Initial planning and requirements gathering"),
    ("Implementation", "Active development and integration"),
    ("Operations", "Production use and maintenance"),
    ("Closeout", "Transition and archival")
]
for phase, desc in phases:
    p = tf.add_paragraph()
    p.text = f"• {phase} → {desc}"
    p.font.size = Pt(18)
    p.space_before = Pt(4)

# ============ SLIDE 10: Demo - Solution Lifecycle ============
slide10 = add_content_slide(prs, "Demo", "Solution Lifecycle Tracking")

# Placeholder for screenshots
placeholder1 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.5), Inches(4))
placeholder1.fill.solid()
placeholder1.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder1.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

placeholder2 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.3), Inches(5.5), Inches(4))
placeholder2.fill.solid()
placeholder2.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder2.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

# LIVE DEMO text
demo_text = slide10.shapes.add_textbox(Inches(4), Inches(5.5), Inches(5), Inches(0.8))
tf = demo_text.text_frame
p = tf.paragraphs[0]
p.text = "LIVE DEMO"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

# Demo instructions
instructions = slide10.shapes.add_textbox(Inches(2), Inches(6.3), Inches(9), Inches(0.8))
tf = instructions.text_frame
p = tf.paragraphs[0]
p.text = "Implementation Dashboard → Select \"Active Solutions\" → MWOW → Details Card"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 11: Q2 - What's the Status? ============
slide11 = add_content_slide(prs, "Q2", "What's the Document & Deliverable Status?")

# Three rounded boxes
add_rounded_box(slide11, Inches(0.3), Inches(1.3), Inches(4.1), Inches(2.3), GREEN, "Document Links", GREEN)
add_rounded_box(slide11, Inches(4.6), Inches(1.3), Inches(4.1), Inches(2.3), BLUE, "Document Status Grid", BLUE)
add_rounded_box(slide11, Inches(8.9), Inches(1.3), Inches(4.1), Inches(2.3), ORANGE, "Milestone Documents", ORANGE)

# Bullets for each box
dl_bullets = slide11.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(3.7), Inches(1.7))
tf = dl_bullets.text_frame
tf.word_wrap = True
for item in ["Drive Folder", "Earthdata Page", "Project Plan", "Science SOW", "Risk Register"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

dsg_bullets = slide11.shapes.add_textbox(Inches(4.8), Inches(1.9), Inches(3.7), Inches(1.7))
tf = dsg_bullets.text_frame
tf.word_wrap = True
for item in ["Project Plan status", "Science SOW status", "IRA/TTA status", "ICD status", "All 9 deliverables tracked"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

md_bullets = slide11.shapes.add_textbox(Inches(9.1), Inches(1.9), Inches(3.7), Inches(1.7))
tf = md_bullets.text_frame
tf.word_wrap = True
for item in ["ATP Memo & Presentation", "F2I Memo & Presentation", "ORR Memo & Presentation", "Closeout Memo & Presentation", "One-click access to all docs"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

# How section
how_label = slide11.shapes.add_textbox(Inches(0.4), Inches(3.8), Inches(12), Inches(0.5))
tf = how_label.text_frame
p = tf.paragraphs[0]
p.text = "How MO Viewer Tracks Document Status:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

how_box = slide11.shapes.add_textbox(Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.5))
tf = how_box.text_frame
tf.word_wrap = True
items = [
    "9 key documents tracked per solution with completion status",
    "Direct links to source documents in Drive",
    "Visual indicators: Complete (green), In Progress (yellow), Not Started (gray)",
    "Quick access to Earthdata page and external resources"
]
for item in items:
    p = tf.add_paragraph()
    p.text = "✓ " + item
    p.font.size = Pt(18)
    p.space_before = Pt(4)

# ============ SLIDE 12: Demo - Document Status ============
slide12 = add_content_slide(prs, "Demo", "Document Status Tracking")

# Placeholder for screenshot
placeholder = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.3), Inches(10), Inches(4))
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

# LIVE DEMO text
demo_text = slide12.shapes.add_textbox(Inches(4), Inches(5.5), Inches(5), Inches(0.8))
tf = demo_text.text_frame
p = tf.paragraphs[0]
p.text = "LIVE DEMO"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

# Demo instructions
instructions = slide12.shapes.add_textbox(Inches(1.5), Inches(6.3), Inches(10), Inches(0.8))
tf = instructions.text_frame
p = tf.paragraphs[0]
p.text = "Implementation Dashboard → HLS → Details Card → Document Status section"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 13: Q3 - What's Happening? ============
slide13 = add_content_slide(prs, "Q3", "What's Happening with This Solution?")

# Three rounded boxes
add_rounded_box(slide13, Inches(0.3), Inches(1.3), Inches(4.1), Inches(2.3), GREEN, "Recent Updates", GREEN)
add_rounded_box(slide13, Inches(4.6), Inches(1.3), Inches(4.1), Inches(2.3), BLUE, "Stakeholder Summary", BLUE)
add_rounded_box(slide13, Inches(8.9), Inches(1.3), Inches(4.1), Inches(2.3), ORANGE, "Team & Contacts", ORANGE)

# Bullets for each box
ru_bullets = slide13.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(3.7), Inches(1.7))
tf = ru_bullets.text_frame
tf.word_wrap = True
for item in ["Updates from meeting notes", "Chronological timeline", "Source document links", "Last 30/60/90 day filters"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

ss_bullets = slide13.shapes.add_textbox(Inches(4.8), Inches(1.9), Inches(3.7), Inches(1.7))
tf = ss_bullets.text_frame
tf.word_wrap = True
for item in ["Total contacts count", "Survey submitters", "Primary/Secondary SMEs", "Agency affiliations"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

tc_bullets = slide13.shapes.add_textbox(Inches(9.1), Inches(1.9), Inches(3.7), Inches(1.7))
tf = tc_bullets.text_frame
tf.word_wrap = True
for item in ["Team lead assignment", "IRA representative", "Agency breakdown", "Email all stakeholders"]:
    add_bullet_paragraph(tf, "• " + item, font_size=14)

# Update tracking section
update_label = slide13.shapes.add_textbox(Inches(0.4), Inches(3.8), Inches(12), Inches(0.5))
tf = update_label.text_frame
p = tf.paragraphs[0]
p.text = "Update Tracking in MO Viewer:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

update_box = slide13.shapes.add_textbox(Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.5))
tf = update_box.text_frame
tf.word_wrap = True
items = [
    "Updates automatically extracted from Internal Planning and SEP meeting notes",
    "Linked to source documents for full context",
    "Filterable by date range and solution",
    "Exportable for reports and presentations"
]
for item in items:
    p = tf.add_paragraph()
    p.text = "→ " + item
    p.font.size = Pt(18)
    p.space_before = Pt(4)

# ============ SLIDE 14: Demo - Updates & Stakeholders ============
slide14 = add_content_slide(prs, "Demo", "Updates and Stakeholder Tracking")

# Placeholder for screenshots
placeholder1 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.5), Inches(4))
placeholder1.fill.solid()
placeholder1.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder1.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

placeholder2 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.3), Inches(5.5), Inches(4))
placeholder2.fill.solid()
placeholder2.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
placeholder2.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)

# LIVE DEMO text
demo_text = slide14.shapes.add_textbox(Inches(4), Inches(5.5), Inches(5), Inches(0.8))
tf = demo_text.text_frame
p = tf.paragraphs[0]
p.text = "LIVE DEMO"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

# Demo instructions
instructions = slide14.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11), Inches(0.8))
tf = instructions.text_frame
p = tf.paragraphs[0]
p.text = "Implementation Dashboard → GABAN → Details Card → Recent Updates / Stakeholders"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 15: Summary ============
slide15 = add_content_slide(prs, "MO Viewer", "Answering Implementation Questions")

# Three summary boxes with colored borders
summary_data = [
    ("Q1: Where Are We?", GREEN, "Solution Picker &\nMilestone Timeline", "track lifecycle phase and\nupcoming milestones"),
    ("Q2: What's the Status?", BLUE, "Document Status Grid\n& Links", "see deliverable completion\nat a glance"),
    ("Q3: What's Happening?", ORANGE, "Recent Updates &\nStakeholders", "stay current on activities\nand contacts")
]

for i, (title, color, feature, desc) in enumerate(summary_data):
    left = Inches(0.4 + i * 4.3)

    # Box with colored border
    box = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(4.1), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(4)

    # Title
    title_box = slide15.shapes.add_textbox(left + Inches(0.2), Inches(1.7), Inches(3.7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Feature
    feat_box = slide15.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.7), Inches(1.2))
    tf = feat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = feature
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide15.shapes.add_textbox(left + Inches(0.2), Inches(4), Inches(3.7), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER

# Tagline
tagline = slide15.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.333), Inches(0.8))
tf = tagline.text_frame
p = tf.paragraphs[0]
p.text = "One platform for complete solution portfolio visibility"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER

# ============ SLIDE 16: Thank You ============
slide16 = prs.slide_layouts[6]
slide16 = prs.slides.add_slide(prs.slide_layouts[6])

# Dark blue background
bg = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Thank You text
thanks = slide16.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = thanks.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Questions text
questions = slide16.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
tf = questions.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(36)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Save the presentation
output_path = r"C:\Users\cjtucke3\Documents\Personal\MO-development\nsite-mo-viewer\MO-Viewer-Implementation-Presentation-v2.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
