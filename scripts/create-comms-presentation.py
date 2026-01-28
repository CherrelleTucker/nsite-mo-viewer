#!/usr/bin/env python3
"""
Create MO Viewer Comms Presentation
Generates a PowerPoint presentation for Comms-NSITE features
Mirrors the style of the SEP presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors (matching SEP presentation)
NAVY = RGBColor(0x1A, 0x23, 0x7E)  # Primary brand color
BLUE = RGBColor(0x30, 0x4F, 0xFE)  # Accent
GREEN = RGBColor(0x4C, 0xAF, 0x50)  # Success/positive
ORANGE = RGBColor(0xFF, 0x98, 0x00)  # Warning/attention
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)  # Comms accent color
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, NAVY)

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "MO Viewer"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Information Management Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Comms-View | January 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: The Problem - Comms specific
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "The Problem: lots of data, not a lot of ways to quickly access the information it informs")

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # MO-wide bullet
    p = tf.paragraphs[0]
    p.text = "MO-wide:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    bullets_mo = [
        "over 9000 interconnected files stored in Google Drive, acting as a \"database\"",
        "Jenny, Cherrelle, Slack, emails, Teams, and various meetings acting as database interfaces"
    ]
    for bullet in bullets_mo:
        p = tf.add_paragraph()
        p.text = "    " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

    # General problem
    p = tf.add_paragraph()
    p.text = "multiple copies of various files used to create actionable information conveyed in meetings → lacking a shared agreement or platform for Source of Truth for information"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(16)

    # Comms-specific
    p = tf.add_paragraph()
    p.text = "Comms-specific:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_before = Pt(20)

    bullets_comms = [
        "~359 Comms-related files (stories, events, outreach, media)",
        "38 stories tracked across multiple spreadsheets and docs",
        "Events scattered across calendars, emails, and planning docs",
        "No single view of solution coverage or messaging gaps",
        "Key messages and blurbs buried in various documents"
    ]
    for bullet in bullets_comms:
        p = tf.add_paragraph()
        p.text = "    " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)

    # Slide 3: Answering Comms Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Answering Comms Questions with CommsViewer")

    questions = [
        ("1", "What Stories?", "What stories are in development and\nwhat's their status?"),
        ("2", "Coverage Gaps?", "Which solutions lack comms coverage\nand need attention?"),
        ("3", "Events & Opportunities?", "What events are coming up and\nhow do we prepare?"),
    ]

    y_start = 1.8
    for i, (num, title_text, desc) in enumerate(questions):
        # Number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_start + i * 1.7), Inches(0.8), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PURPLE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

        # Title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + i * 1.7), Inches(4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PURPLE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + 0.5 + i * 1.7), Inches(10), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY

    # Slide 4: Q1 - What Stories
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q1: What Stories Are In Development?", PURPLE)

    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2.5),
        "Story Pipeline",
        [
            "Kanban-style story tracking",
            "Status: Draft → Review → Published",
            "Solution linkage for each story",
            "Content type categorization"
        ],
        PURPLE)

    add_feature_box(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(2.5),
        "Story Opportunities",
        [
            "Auto-detected from milestones",
            "ATP, ORR, F2I triggers",
            "One-click story creation",
            "Links to source updates"
        ],
        GREEN)

    add_feature_box(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(2.5),
        "Key Messages",
        [
            "Solution-specific messaging",
            "Highlighter blurbs library",
            "Searchable message bank",
            "Priority alignment tags"
        ],
        BLUE)

    # Answer box
    answer = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.5))
    tf = answer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Story Tracking in MO Viewer:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    bullets = [
        "All stories visible in pipeline view with drag-and-drop status updates",
        "Story opportunities auto-generated from solution milestones",
        "Key messages searchable and linked to solutions",
        "Admin priorities alignment (Partnerships, AI, Science Integrity, etc.)"
    ]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Slide 5: Demo - Stories
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Story Pipeline & Opportunities", PURPLE)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = PURPLE

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nComms → Pipeline → Story Details → Key Messages"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 6: Q2 - Coverage Gaps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q2: Which Solutions Lack Coverage?", PURPLE)

    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
        "Coverage Analysis",
        [
            "Visual coverage map by solution",
            "Gap identification & alerts",
            "Story count per solution",
            "Last coverage date tracking"
        ],
        ORANGE)

    add_feature_box(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
        "Priority Alignment",
        [
            "Admin priorities dashboard",
            "Biden-Harris alignment tags",
            "Partnership opportunities",
            "Science advancement tracking"
        ],
        GREEN)

    # Metrics
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.8))
    tf = metrics_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "How MO Viewer Identifies Coverage Gaps:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    points = [
        "Solutions without recent stories flagged automatically",
        "Coverage gaps panel shows solutions needing attention",
        "Filters by lifecycle phase, cycle, and content type",
        "One-click navigation to create new story for gap"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Slide 7: Demo - Coverage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Coverage Analysis", PURPLE)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = PURPLE

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nComms → Coverage → Gaps Panel → Priorities View"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 8: Q3 - Events & Opportunities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q3: Events & Outreach Opportunities?", PURPLE)

    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2.5),
        "Events Pipeline",
        [
            "Track conferences & events",
            "Status: Potential → Confirmed",
            "Guest list management",
            "Deadline tracking"
        ],
        BLUE)

    add_feature_box(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(2.5),
        "Event Prep Reports",
        [
            "Auto-generated briefings",
            "Guest profiles & agencies",
            "Conversation starters",
            "Export to Google Doc"
        ],
        GREEN)

    add_feature_box(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(2.5),
        "Calendar View",
        [
            "Visual event timeline",
            "Sector-based filtering",
            "Upcoming deadlines",
            "Team coordination"
        ],
        ORANGE)

    # Events section
    events_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.8))
    tf = events_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Event Management in MO Viewer:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    points = [
        "Track all outreach events from potential to attended",
        "Build guest lists with stakeholder connections",
        "Generate prep reports with talking points and context",
        "Export briefing docs for meetings and travel"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = "→ " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Slide 9: Demo - Events
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Event Prep & Guest Management", PURPLE)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = PURPLE

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nComms → Events → Guest List → Prep Report → Export to Doc"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 10: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "MO Viewer: Answering Comms Questions", PURPLE)

    summaries = [
        ("Q1: What Stories?", "Pipeline view tracks all stories\nfrom draft to published with\nopportunity detection", GREEN),
        ("Q2: Coverage Gaps?", "Coverage analysis identifies\nsolutions needing attention\nwith priority alignment", BLUE),
        ("Q3: Events?", "Event management with\nguest lists, prep reports,\nand doc export", ORANGE),
    ]

    x_positions = [0.5, 4.5, 8.5]
    for i, (title_text, desc, color) in enumerate(summaries):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.8), Inches(4), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # Title inside box
        title_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2), Inches(3.6), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2.8), Inches(3.6), Inches(2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Bottom tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "One platform for complete communications visibility"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Slide 11: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    questions_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = questions_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), '..', 'MO-Viewer-Comms-Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

def add_background(slide, color):
    """Add a solid background color to slide"""
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_header(slide, text, color=None):
    """Add a header to the slide"""
    if color is None:
        color = NAVY
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

def add_feature_box(slide, left, top, width, height, title, bullets, accent_color):
    """Add a feature box with title and bullets"""
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent_color
    box.line.width = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_color

    # Bullets
    bullets_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)

if __name__ == "__main__":
    create_presentation()
