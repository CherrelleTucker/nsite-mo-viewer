# -*- coding: utf-8 -*-
"""Quick exploration of PowerPoint structure"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

pptx_path = Path(r"C:\Users\cjtucke3\Documents\Personal\MO-development\source-archives\Monthly Project Status Updates\2026-01 NSITE Monthly Meeting for Solution Updates Presentation.pptx")

prs = Presentation(pptx_path)

print(f"File: {pptx_path.name}")
print(f"Slides: {len(prs.slides)}")
print()

# Look at slides 7-15 (likely solution update slides)
for i, slide in enumerate(prs.slides):
    if i < 7 or i > 15:
        continue

    print(f"=== Slide {i+1} ===")

    # Get slide title if present
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text[:100]}")

    # Iterate through shapes and show full text
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            shape_type = type(shape).__name__
            # Print more of each shape to see content structure
            print(f"  [{shape_type}]:")
            for line in text.split('\n')[:10]:
                if line.strip():
                    print(f"    {line.strip()[:120]}")
            if len(text.split('\n')) > 10:
                print(f"    ... ({len(text.split(chr(10)))} lines total)")

    print()
