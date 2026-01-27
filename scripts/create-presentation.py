#!/usr/bin/env python3
"""
Create MO Viewer SEP Presentation
Generates a PowerPoint presentation answering the 3 SEP questions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
NAVY = RGBColor(0x1A, 0x23, 0x7E)  # Primary brand color
BLUE = RGBColor(0x30, 0x4F, 0xFE)  # Accent
GREEN = RGBColor(0x4C, 0xAF, 0x50)  # Success/positive
ORANGE = RGBColor(0xFF, 0x98, 0x00)  # Warning/attention
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, NAVY)

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "MO Viewer"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Outreach Stakeholder Engagement Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SEP Review | January 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Agenda / 3 Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Today's SEP Questions")

    questions = [
        ("1", "What's New?", "What is the most recent thing this Solution did\nand did we talk to stakeholders about it?"),
        ("2", "Meeting Needs?", "Are we meeting the needs of stakeholders\nwith this solution?"),
        ("3", "Growth Opportunities?", "Are there unexplored agencies we can\nconnect with to promote this Solution?"),
    ]

    y_start = 1.8
    for i, (num, title, desc) in enumerate(questions):
        # Number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_start + i * 1.7), Inches(0.8), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

        # Title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + i * 1.7), Inches(4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Description
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + 0.5 + i * 1.7), Inches(10), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY

    # Slide 3: Question 1 - What's New
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q1: What's New & Did We Discuss It?", NAVY)

    # Implementation Dashboard highlight
    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
        "Implementation Dashboard",
        [
            "Tracks all SEP activities with status",
            "Color-coded progress indicators",
            "Stakeholder touchpoints logged",
            "Historical engagement timeline"
        ],
        GREEN)

    add_feature_box(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
        "SEP Dashboard",
        [
            "Real-time stakeholder metrics",
            "Communication logs & history",
            "Meeting notes & follow-ups",
            "Engagement frequency tracking"
        ],
        BLUE)

    # Answer box
    answer = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.5))
    tf = answer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Recent Activities Tracked in MO Viewer:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY

    bullets = [
        "All stakeholder communications logged with dates and outcomes",
        "Event participation tracked (conferences, webinars, meetings)",
        "Reports generated and shared with leadership",
        "Follow-up actions assigned and monitored"
    ]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Slide 4: Question 1 - Demo View
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Implementation & Activity Tracking", NAVY)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = NAVY

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nImplementation Dashboard → Activity Timeline → Reports"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 5: Question 2 - Meeting Needs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q2: Are We Meeting Stakeholder Needs?", NAVY)

    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2.2),
        "Agencies View",
        [
            "Complete agency profiles",
            "Contact information",
            "Engagement history",
            "Need alignment status"
        ],
        BLUE)

    add_feature_box(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(2.2),
        "Contacts Directory",
        [
            "Stakeholder database",
            "Relationship tracking",
            "Communication prefs",
            "Role-based filtering"
        ],
        GREEN)

    add_feature_box(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(2.2),
        "Need Alignment Report",
        [
            "Gap analysis",
            "Coverage metrics",
            "Priority mapping",
            "Action items"
        ],
        ORANGE)

    # Key metrics
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(3))
    tf = metrics_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "How MO Viewer Ensures We Meet Needs:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY

    points = [
        "Stakeholder needs mapped to solution capabilities",
        "Engagement scoring identifies under-served agencies",
        "Automated reports flag gaps in coverage",
        "Historical data shows improvement trends"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = "✓ " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(10)

    # Slide 6: Question 2 - Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Stakeholder Management", NAVY)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = NAVY

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nSEP Dashboard → Agencies → Contacts → Need Alignment Report"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 7: Question 3 - Growth Opportunities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Q3: Unexplored Growth Opportunities?", NAVY)

    add_feature_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
        "Cold Agency Indicators",
        [
            "Agencies with no recent contact",
            "Engagement score below threshold",
            "Visual indicators (blue = cold)",
            "Prioritization recommendations"
        ],
        BLUE)

    add_feature_box(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5),
        "Department Reach Report",
        [
            "Coverage by department",
            "Penetration percentages",
            "White space analysis",
            "Target recommendations"
        ],
        ORANGE)

    # Outreach section
    outreach_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.8))
    tf = outreach_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Outreach & Events Module"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY

    points = [
        "Track upcoming conferences and events",
        "Identify networking opportunities",
        "Generate prep reports with guest profiles",
        "Export to Google Docs for meetings"
    ]
    for point in points:
        p = tf.add_paragraph()
        p.text = "→ " + point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Slide 8: Question 3 - Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Demo: Finding Growth Opportunities", NAVY)

    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = LIGHT_GRAY
    demo_box.line.color.rgb = NAVY

    tf = demo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\n\nLIVE DEMO"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nAgencies (Cold) → Dept Reach Report → Events → Prep Report"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Slide 9: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "MO Viewer: Answering SEP Questions", NAVY)

    summaries = [
        ("Q1: What's New?", "Implementation Dashboard & Activity Logs\ntrack all recent stakeholder interactions", GREEN),
        ("Q2: Meeting Needs?", "SEP Dashboard, Agencies & Contacts views\nensure comprehensive stakeholder coverage", BLUE),
        ("Q3: Growth?", "Cold indicators, Dept Reach & Events\nidentify and pursue new opportunities", ORANGE),
    ]

    x_positions = [0.5, 4.5, 8.5]
    for i, (title, desc, color) in enumerate(summaries):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_positions[i]), Inches(1.8), Inches(4), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(3)

        # Title inside box
        title_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2), Inches(3.6), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(2.8), Inches(3.6), Inches(2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

    # Bottom tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(1))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "One platform for complete stakeholder engagement visibility"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    questions_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = questions_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), '..', 'MO-Viewer-SEP-Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

def add_background(slide, color):
    """Add a solid background color to slide"""
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_header(slide, text, color=None):
    """Add a header to the slide"""
    if color is None:
        color = NAVY
    header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color

def add_feature_box(slide, left, top, width, height, title, bullets, accent_color):
    """Add a feature box with title and bullets"""
    # Box background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent_color
    box.line.width = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_color

    # Bullets
    bullets_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)

if __name__ == "__main__":
    create_presentation()
